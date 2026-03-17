
import pymupdf as fitz
from pptx import Presentation
from docx import Document
import os


def create_txt_file(file_path: str, text: str, destination: str):
    name = os.path.splitext(file_path)[0].split('\\')[-1]
    print(f" Creating: {name}.txt/n path: {file_path}")
    path = os.path.join(destination, f"{name}.txt")
    with open(path, "w", encoding="utf-8") as f:
        f.write(text)
    print(f"✅ Created: {path}")

def pdf_to_text(file_path: str) -> str:
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def docx_to_text(file_path: str) -> str:
    try:
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        print(f"❌ Error reading docx: {file_path} → {e}")
        return ""


def pptx_to_text(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + "\n"
    return text

def extract_text(file_path: str) -> str:
    ext = file_path.split(".")[-1].lower()
    
    if ext == "pdf":
        return pdf_to_text(file_path)
    elif ext == "docx":
        return docx_to_text(file_path)
    elif ext == "pptx":
        return pptx_to_text(file_path)
    elif ext == "txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    else:
        print(f"Unsupported file type: {ext}")
        return ""
    
async def get_txt_data(folder: str = "user_files",destination: str = "data"):
    for filename in os.listdir(folder):
        full_path = os.path.join(folder, filename)
        if os.path.isfile(full_path):
            create_txt_file(full_path,extract_text(full_path),destination)


        
